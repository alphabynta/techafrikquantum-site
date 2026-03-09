#!/usr/bin/env python3
"""Generate TechAfrik business card PPTX — Madjou Balde (dark + white versions)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ────────────────────────────────────────────────────────────────────
AMBER      = RGBColor(0xF5, 0xA5, 0x24)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
NEAR_WHITE = RGBColor(0xF0, 0xF0, 0xF2)
CONTACT_DK = RGBColor(0x55, 0x55, 0x55)   # contact text on dark card
CONTACT_LT = RGBColor(0x99, 0x99, 0x99)   # contact text on white card

# ── Logo line segments  (SVG viewBox 0 0 116 44) ─────────────────────────────
# kind: 'box'=amber border, 'tech'=white/black letters, 'afrik'=amber letters
LOGO_LINES = [
    # TECH amber box
    ( 0,16, 51,16,'box'), ( 0,36, 51,36,'box'), ( 0,16,  0,36,'box'), (51,16, 51,36,'box'),
    # T
    ( 3,19, 13,19,'tech'), ( 8,19,  8,34,'tech'),
    # E
    (25,19, 15,19,'tech'), (15,19, 15,34,'tech'), (15,34, 25,34,'tech'), (15,26, 23,26,'tech'),
    # C
    (37,19, 27,19,'tech'), (27,19, 27,34,'tech'), (27,34, 37,34,'tech'),
    # H
    (39,19, 39,34,'tech'), (49,19, 49,34,'tech'), (39,26, 49,26,'tech'),
    # A
    (55,34, 60,19,'afrik'), (60,19, 65,34,'afrik'), (57,27, 63,27,'afrik'),
    # F
    (77,19, 67,19,'afrik'), (67,19, 67,34,'afrik'), (67,26, 75,26,'afrik'),
    # R
    (79,19, 79,34,'afrik'), (79,19, 89,19,'afrik'), (89,19, 89,26,'afrik'),
    (89,26, 79,26,'afrik'), (83,26, 89,34,'afrik'),
    # I
    (96,19, 96,34,'afrik'), (91,19,101,19,'afrik'), (91,34,101,34,'afrik'),
    # K
    (103,19,103,34,'afrik'), (103,26,113,19,'afrik'), (103,26,113,34,'afrik'),
]

def draw_logo(slide, x_in, y_in, width_in, tech_color=WHITE, stroke_pt=1.4):
    """Render the TechAfrik logo at (x_in, y_in) scaled to width_in."""
    scale = Inches(width_in) / 116       # EMU per SVG unit
    ox, oy = Inches(x_in), Inches(y_in)
    for x1, y1, x2, y2, kind in LOGO_LINES:
        color = AMBER if kind in ('box', 'afrik') else tech_color
        ex1 = int(ox + x1 * scale)
        ey1 = int(oy + y1 * scale)
        ex2 = int(ox + x2 * scale)
        ey2 = int(oy + y2 * scale)
        conn = slide.shapes.add_connector(1, ex1, ey1, ex2, ey2)
        conn.line.color.rgb = color
        conn.line.width = Pt(stroke_pt)

def add_txt(slide, text, l, t, w, h, font_sz, bold, color, align=PP_ALIGN.LEFT,
            letter_spacing=None, caps=False):
    """Add a single-line text box."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'Inter'
    run.font.size = Pt(font_sz)
    run.font.bold = bold
    run.font.color.rgb = color
    if caps:
        from pptx.oxml.ns import qn
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('w:caps'), '1')

def add_connector(slide, x1_in, y1_in, x2_in, y2_in, color, width_pt):
    conn = slide.shapes.add_connector(
        1,
        Inches(x1_in), Inches(y1_in),
        Inches(x2_in), Inches(y2_in),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Presentation setup: 4 slides at business-card size ───────────────────────
prs = Presentation()
prs.slide_width  = Inches(3.5)
prs.slide_height = Inches(2.0)
blank = prs.slide_layouts[6]  # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# Shared layout constants
LOGO_X      = 0.22   # logo left margin
LOGO_Y      = 0.20   # logo top
LOGO_W      = 1.20   # logo width

# Left column (name + title)
NAME_X      = 0.22
NAME_Y      = 0.87
NAME_W      = 1.65

TITLE_X     = 0.22
TITLE_Y     = 1.27
TITLE_W     = 1.65

RULE_X1     = 0.22
RULE_Y      = 1.57
RULE_X2     = 0.50

# Right column (contacts) — aligned to the right portion of card
CONTACT_X   = 2.05
PHONE_Y     = 1.00
EMAIL_Y     = 1.30
CONTACT_W   = 1.28

# Thin vertical amber divider
DIVIDER_X   = 1.90
DIVIDER_Y1  = 0.85
DIVIDER_Y2  = 1.70

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Dark Front
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
set_bg(s1, BLACK)
draw_logo(s1, LOGO_X, LOGO_Y, LOGO_W, tech_color=WHITE)

add_txt(s1, 'Madjou Balde',   NAME_X, NAME_Y,  NAME_W, 0.42, 15.5, True,  NEAR_WHITE)
add_txt(s1, 'PRESIDENT & CEO', TITLE_X, TITLE_Y, TITLE_W, 0.20, 6.0, True, AMBER)
add_connector(s1, RULE_X1, RULE_Y, RULE_X2, RULE_Y, RGBColor(0xF5, 0xA5, 0x24), 0.75)

# Vertical divider
add_connector(s1, DIVIDER_X, DIVIDER_Y1, DIVIDER_X, DIVIDER_Y2, RGBColor(0x2A,0x2A,0x2A), 0.5)

add_txt(s1, '+224 xxx yy uu ww',           CONTACT_X, PHONE_Y, CONTACT_W, 0.20, 6.5, False, CONTACT_DK)
add_txt(s1, 'madjou.balde@techafrik.com',  CONTACT_X, EMAIL_Y, CONTACT_W, 0.20, 6.5, False, CONTACT_DK)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Dark Back
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
set_bg(s2, BLACK)
draw_logo(s2, 0.72, 0.62, 2.06, tech_color=WHITE, stroke_pt=1.8)   # large centered logo

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — White Front
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
set_bg(s3, WHITE)
draw_logo(s3, LOGO_X, LOGO_Y, LOGO_W, tech_color=BLACK)

add_txt(s3, 'Madjou Balde',    NAME_X, NAME_Y,  NAME_W, 0.42, 15.5, True,  BLACK)
add_txt(s3, 'PRESIDENT & CEO', TITLE_X, TITLE_Y, TITLE_W, 0.20, 6.0, True, AMBER)
add_connector(s3, RULE_X1, RULE_Y, RULE_X2, RULE_Y, AMBER, 0.75)

add_connector(s3, DIVIDER_X, DIVIDER_Y1, DIVIDER_X, DIVIDER_Y2, RGBColor(0xD4,0xD4,0xD4), 0.5)

add_txt(s3, '+224 xxx yy uu ww',           CONTACT_X, PHONE_Y, CONTACT_W, 0.20, 6.5, False, CONTACT_LT)
add_txt(s3, 'madjou.balde@techafrik.com',  CONTACT_X, EMAIL_Y, CONTACT_W, 0.20, 6.5, False, CONTACT_LT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — White Back
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank)
set_bg(s4, WHITE)
draw_logo(s4, 0.72, 0.62, 2.06, tech_color=BLACK, stroke_pt=1.8)

# ─────────────────────────────────────────────────────────────────────────────
out = '/home/samuel/alimou/website/business-card-madjou-balde.pptx'
prs.save(out)
print(f'Saved: {out}')
print('Slides: 1=Dark Front, 2=Dark Back, 3=White Front, 4=White Back')
